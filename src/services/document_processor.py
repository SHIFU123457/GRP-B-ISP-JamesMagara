import os
import logging
from pathlib import Path
from typing import List, Dict, Any, Optional
import PyPDF2
import fitz  # PyMuPDF for better PDF extraction
from docx import Document as DocxDocument
from pptx import Presentation  # PowerPoint processing
import nltk
from nltk.tokenize import sent_tokenize, word_tokenize
from nltk.corpus import stopwords

# OCR imports (optional - will gracefully degrade if not available)
try:
    from pdf2image import convert_from_path
    import pytesseract
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False
    logging.warning("OCR libraries (pdf2image, pytesseract) not installed. Scanned PDFs will not be processed.")

logger = logging.getLogger(__name__)

def safe_file_operation(operation_name: str):
    """Decorator for safe file operations with detailed error reporting"""
    def decorator(func):
        def wrapper(*args, **kwargs):
            try:
                return func(*args, **kwargs)
            except FileNotFoundError as e:
                logger.error(f"{operation_name} failed - File not found: {e}")
                return None
            except PermissionError as e:
                logger.error(f"{operation_name} failed - Permission denied: {e}")
                return None
            except Exception as e:
                logger.error(f"{operation_name} failed - Unexpected error: {e}", exc_info=True)
                return None
        return wrapper
    return decorator

class DocumentProcessor:
    """Process documents and extract text content for RAG pipeline"""
    
    def __init__(self):
        from config.settings import settings
        self.chunk_size = settings.CHUNK_SIZE  # words per chunk (from config)
        self.chunk_overlap = settings.CHUNK_OVERLAP  # overlap between chunks
        
        # Download NLTK data if not present
        try:
            nltk.data.find('tokenizers/punkt')
            nltk.data.find('corpora/stopwords')
        except LookupError:
            nltk.download('punkt', quiet=True)
            nltk.download('stopwords', quiet=True)
    
    @safe_file_operation("Text extraction")
    def extract_text(self, file_path: str, file_type: str) -> Optional[str]:
        """Extract text from various document formats"""
        try:
            if file_type.lower() == 'pdf':
                return self._extract_pdf_text(file_path)
            elif file_type.lower() in ['docx', 'doc']:
                return self._extract_docx_text(file_path)
            elif file_type.lower() in ['pptx', 'ppt']:
                return self._extract_pptx_text(file_path)
            elif file_type.lower() == 'txt':
                return self._extract_txt_text(file_path)
            else:
                logger.warning(f"Unsupported file type: {file_type}")
                return None

        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {e}")
            return None
    
    @safe_file_operation("PDF text extraction")
    def _extract_pdf_text(self, file_path: str) -> str:
        """Extract text from PDF files using PyMuPDF (fallback to PyPDF2, then OCR)"""
        text = ""

        try:
            # Try PyMuPDF first (better extraction)
            doc = fitz.open(file_path)
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text += page.get_text() + "\n"
            doc.close()

            # If we got good text, return it
            if len(text.strip()) > 50:  # Reasonable amount of text
                logger.info(f"Successfully extracted {len(text)} characters using PyMuPDF")
                return text.strip()
        except Exception as e:
            logger.warning(f"PyMuPDF extraction failed for {file_path}: {e}, trying PyPDF2...")

        # Fallback to PyPDF2
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"

            # Check if PyPDF2 got good text
            if len(text.strip()) > 50:
                logger.info(f"Extracted {len(text)} characters using PyPDF2 fallback")
                return text.strip()
            else:
                logger.warning(f"PyPDF2 extracted insufficient text ({len(text)} chars), trying OCR...")
        except Exception as e:
            logger.warning(f"PyPDF2 extraction failed for {file_path}: {e}, trying OCR...")

        # Final fallback: OCR for scanned/image PDFs
        if OCR_AVAILABLE:
            try:
                ocr_text = self._extract_pdf_with_ocr(file_path)
                if ocr_text and len(ocr_text.strip()) > 50:
                    logger.info(f"Successfully extracted {len(ocr_text)} characters using OCR")
                    return ocr_text.strip()
                else:
                    logger.error(f"OCR extraction failed or returned insufficient text from {file_path}")
            except Exception as e:
                logger.error(f"OCR extraction failed for {file_path}: {e}")
        else:
            logger.error(f"OCR not available. Cannot process scanned PDF: {file_path}")
            logger.info("Install OCR support with: pip install pdf2image pytesseract")
            logger.info("Also install Tesseract OCR: https://github.com/tesseract-ocr/tesseract")

        return text.strip()

    def _extract_pdf_with_ocr(self, file_path: str) -> str:
        """Extract text from PDF using OCR (for scanned/image PDFs)"""
        if not OCR_AVAILABLE:
            logger.error("OCR libraries not available")
            return ""

        try:
            logger.info(f"Converting PDF to images for OCR: {file_path}")
            # Convert PDF pages to images
            images = convert_from_path(file_path, dpi=300)  # Higher DPI for better OCR accuracy

            logger.info(f"Running OCR on {len(images)} pages...")
            text_parts = []

            for i, image in enumerate(images, 1):
                logger.info(f"Processing page {i}/{len(images)} with OCR...")
                # Extract text from image using Tesseract
                page_text = pytesseract.image_to_string(image, lang='eng')

                if page_text.strip():
                    text_parts.append(f"--- Page {i} ---\n{page_text}")
                    logger.debug(f"Page {i}: Extracted {len(page_text)} characters")
                else:
                    logger.warning(f"Page {i}: No text extracted")

            final_text = "\n\n".join(text_parts)
            logger.info(f"OCR completed: {len(final_text)} total characters from {len(images)} pages")

            return final_text

        except Exception as e:
            logger.error(f"OCR processing failed: {e}")
            # Check if Tesseract is installed
            try:
                import subprocess
                subprocess.run(['tesseract', '--version'], capture_output=True, check=True)
            except (subprocess.CalledProcessError, FileNotFoundError):
                logger.error("Tesseract OCR engine not found. Please install it:")
                logger.error("  Windows: Download from https://github.com/UB-Mannheim/tesseract/wiki")
                logger.error("  Linux: sudo apt-get install tesseract-ocr")
                logger.error("  macOS: brew install tesseract")
            return ""

    @safe_file_operation("DOCX text extraction")
    def _extract_docx_text(self, file_path: str) -> str:
        """Extract text from Word documents"""
        doc = DocxDocument(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text.strip()
    
    @safe_file_operation("TXT text extraction")
    def _extract_txt_text(self, file_path: str) -> str:
        """Extract text from plain text files"""
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read().strip()

    @safe_file_operation("PowerPoint text extraction")
    def _extract_pptx_text(self, file_path: str) -> str:
        """Extract text from PowerPoint files (.pptx, .ppt)"""
        try:
            presentation = Presentation(file_path)
            text_content = []

            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text = []

                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())

                    # Handle tables in slides
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                slide_text.append(" | ".join(row_text))

                # Add slide content with slide number
                if slide_text:
                    slide_content = f"\n--- Slide {slide_num} ---\n" + "\n".join(slide_text)
                    text_content.append(slide_content)

            extracted_text = "\n\n".join(text_content)
            logger.info(f"Successfully extracted {len(extracted_text)} characters from PowerPoint with {len(presentation.slides)} slides")

            return extracted_text.strip()

        except Exception as e:
            logger.error(f"Error extracting text from PowerPoint {file_path}: {e}")
            return ""
    
    def clean_text(self, text: str) -> str:
        """Clean and preprocess text"""
        # Remove excessive whitespace
        text = ' '.join(text.split())
        
        # Remove special characters but keep periods and commas
        import re
        text = re.sub(r'[^\w\s.,!?-]', ' ', text)
        
        # Remove multiple spaces
        text = re.sub(r'\s+', ' ', text).strip()
        
        return text
    
    def chunk_text(self, text: str, metadata: Dict[str, Any]) -> List[Dict[str, Any]]:
        """Split text into overlapping chunks for better retrieval"""
        # Clean the text first
        text = self.clean_text(text)
        
        # Split into sentences
        sentences = sent_tokenize(text)
        
        chunks = []
        current_chunk = []
        current_word_count = 0
        
        for sentence in sentences:
            words = word_tokenize(sentence)
            sentence_word_count = len(words)
            
            # If adding this sentence exceeds chunk size, save current chunk
            if current_word_count + sentence_word_count > self.chunk_size and current_chunk:
                chunk_text = ' '.join(current_chunk)
                chunks.append(self._create_chunk(chunk_text, metadata, len(chunks)))
                
                # Start new chunk with overlap
                overlap_sentences = current_chunk[-2:] if len(current_chunk) >= 2 else current_chunk
                current_chunk = overlap_sentences + [sentence]
                current_word_count = len(word_tokenize(' '.join(current_chunk)))
            else:
                current_chunk.append(sentence)
                current_word_count += sentence_word_count
        
        # Add the last chunk if it has content
        if current_chunk:
            chunk_text = ' '.join(current_chunk)
            chunks.append(self._create_chunk(chunk_text, metadata, len(chunks)))
        
        return chunks
    
    def _create_chunk(self, text: str, metadata: Dict[str, Any], chunk_index: int) -> Dict[str, Any]:
        """Create a chunk object with text and metadata"""
        return {
            'text': text,
            'metadata': {
                **metadata,
                'chunk_index': chunk_index,
                'chunk_length': len(text.split()),
            }
        }
    
    def process_document(self, file_path: str, file_type: str, metadata: Dict[str, Any]) -> List[Dict[str, Any]]:
        """Complete document processing pipeline"""
        logger.info(f"Processing document: {file_path}")
        
        # Extract text
        text = self.extract_text(file_path, file_type)
        if not text:
            logger.error(f"No text extracted from {file_path}")
            return []
        
        # Create chunks
        chunks = self.chunk_text(text, metadata)
        
        logger.info(f"Created {len(chunks)} chunks from {file_path}")
        return chunks